"""File content extraction registry — handles txt, docx, pptx, pdf, xlsx, csv, eml."""

from pathlib import Path
from typing import Callable

from core.logging import log


def _extract_plaintext(filepath: Path) -> str | None:
    """Extract text from plain text files (.txt, .md, .vtt, .csv, .eml)."""
    try:
        return filepath.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return filepath.read_text(encoding="latin-1")


def _extract_docx(filepath: Path) -> str | None:
    """Extract text from Word documents (.docx)."""
    import docx
    doc = docx.Document(str(filepath))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(filepath: Path) -> str | None:
    """Extract text from PowerPoint files (.pptx)."""
    from pptx import Presentation
    prs = Presentation(str(filepath))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [shape.text for shape in slide.shapes
                      if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def _extract_pdf(filepath: Path) -> str | None:
    """Extract text from PDF files (.pdf)."""
    import PyPDF2
    with open(filepath, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                text_parts.append(text)
        return "\n\n".join(text_parts)


def _extract_xlsx(filepath: Path) -> str | None:
    """Extract text from Excel files (.xlsx)."""
    import openpyxl
    wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows[:200]))
    wb.close()
    return "\n\n".join(text_parts)


# Map file extensions to extractor functions
EXTRACTORS: dict[str, Callable] = {
    ".txt": _extract_plaintext,
    ".md": _extract_plaintext,
    ".vtt": _extract_plaintext,
    ".csv": _extract_plaintext,
    ".eml": _extract_plaintext,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
    ".xlsx": _extract_xlsx,
}


def extract_text(filepath: Path) -> str | None:
    """Extract text content from a file using the registered extractor.

    Returns text content or None if the file type isn't supported or extraction fails.
    """
    ext = filepath.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return None

    try:
        return extractor(filepath)
    except ImportError as e:
        log.warning(f"    Missing dependency for {ext}: {e} — pip install it to enable")
        return None
    except Exception as e:
        log.warning(f"    Failed to read {filepath.name}: {e}")
        return None
