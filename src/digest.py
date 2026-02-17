"""Internal digest — reads local content and generates a structured daily digest.

Two-phase approach:
1. Collection: Scan input folders for new/modified files, extract text content
2. Analysis: Send content to GHCP SDK agent for structured summarization

The agent generates: TLDRs, decisions, action items, risk flags, and a
"Needs Your Attention" section — all based on standing instructions.
"""

import json
from datetime import datetime
from pathlib import Path

from copilot import CopilotClient

from intel import collect_feeds
from session import PROJECT_ROOT, OUTPUT_DIR, _load_instruction
from tools import get_tools, _load_actions
from utils import agent_session, log


# Max characters to send per file (avoid blowing up the context window)
MAX_CHARS_PER_FILE = 50_000
# Max total characters for all content in a single digest run
MAX_TOTAL_CHARS = 400_000


# --- Phase 1: Content Collection ---

def _load_digest_state(state_file: Path) -> dict:
    """Load the incremental processing state (which files have been digested)."""
    if state_file.exists():
        return json.loads(state_file.read_text(encoding="utf-8"))
    return {"processed": {}}


def _save_digest_state(state_file: Path, state: dict):
    """Save the incremental processing state."""
    state_file.parent.mkdir(parents=True, exist_ok=True)
    state_file.write_text(json.dumps(state, indent=2), encoding="utf-8")


# ---------------------------------------------------------------------------
# File text extractors — registry pattern for clean extensibility
# ---------------------------------------------------------------------------

def _extract_plaintext(filepath: Path) -> str | None:
    """Extract text from plain text files (.txt, .md, .vtt, .csv, .eml)."""
    try:
        return filepath.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return filepath.read_text(encoding="latin-1")


def _extract_docx(filepath: Path) -> str | None:
    """Extract text from Word documents (.docx)."""
    import docx
    doc = docx.Document(str(filepath))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(filepath: Path) -> str | None:
    """Extract text from PowerPoint files (.pptx)."""
    from pptx import Presentation
    prs = Presentation(str(filepath))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [shape.text for shape in slide.shapes
                      if hasattr(shape, "text") and shape.text.strip()]
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def _extract_pdf(filepath: Path) -> str | None:
    """Extract text from PDF files (.pdf)."""
    import PyPDF2
    with open(filepath, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        text_parts = [page.extract_text() for page in reader.pages
                      if page.extract_text() and page.extract_text().strip()]
        return "\n\n".join(text_parts)


def _extract_xlsx(filepath: Path) -> str | None:
    """Extract text from Excel files (.xlsx)."""
    import openpyxl
    wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows[:200]))
    wb.close()
    return "\n\n".join(text_parts)


# Map file extensions to extractor functions
EXTRACTORS: dict[str, callable] = {
    ".txt": _extract_plaintext,
    ".md": _extract_plaintext,
    ".vtt": _extract_plaintext,
    ".csv": _extract_plaintext,
    ".eml": _extract_plaintext,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
    ".xlsx": _extract_xlsx,
}


def _extract_text(filepath: Path) -> str | None:
    """Extract text content from a file using the registered extractor.

    Returns text content or None if the file type isn't supported or extraction fails.
    """
    ext = filepath.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        return None

    try:
        return extractor(filepath)
    except ImportError as e:
        # Missing optional dependency — warn once and skip
        log.warning(f"    Missing dependency for {ext}: {e} — pip install it to enable")
        return None
    except Exception as e:
        log.warning(f"    Failed to read {filepath.name}: {e}")
        return None


def collect_content(config: dict) -> list[dict]:
    """Scan configured input folders and collect text content from new files.

    Returns a list of dicts: {path, type, name, content, size}
    """
    digest_cfg = config.get("digest", {})
    input_paths = digest_cfg.get("input_paths", [])
    supported_ext = set(digest_cfg.get("supported_extensions",
                                        [".vtt", ".txt", ".md", ".docx", ".pptx",
                                         ".pdf", ".xlsx", ".csv", ".eml"]))
    incremental = digest_cfg.get("incremental", True)
    state_file = PROJECT_ROOT / digest_cfg.get("state_file", "output/.digest-state.json")

    state = _load_digest_state(state_file) if incremental else {"processed": {}}
    processed = state.get("processed", {})
    collected = []
    total_chars = 0

    for path_cfg in input_paths:
        folder = Path(path_cfg["path"])
        # Support both absolute and relative (to project root) paths
        if not folder.is_absolute():
            folder = PROJECT_ROOT / folder
        content_type = path_cfg.get("type", "unknown")

        if not folder.exists():
            log.info(f"  Input path does not exist (creating): {folder}")
            folder.mkdir(parents=True, exist_ok=True)
            continue

        log.info(f"  Scanning {folder} (type: {content_type})...")

        for filepath in sorted(folder.rglob("*")):
            if not filepath.is_file():
                continue
            if filepath.suffix.lower() not in supported_ext:
                continue
            if filepath.name.startswith("."):
                continue

            # Incremental: skip files already processed (same path + same mtime)
            file_key = str(filepath)
            file_mtime = filepath.stat().st_mtime
            if incremental and file_key in processed:
                if processed[file_key] >= file_mtime:
                    continue

            # Extract text
            text = _extract_text(filepath)
            if not text or not text.strip():
                log.debug(f"    SKIP (empty/unreadable): {filepath.name}")
                continue

            # Truncate if too large
            if len(text) > MAX_CHARS_PER_FILE:
                text = text[:MAX_CHARS_PER_FILE] + f"\n\n[... truncated at {MAX_CHARS_PER_FILE} chars]"

            if total_chars + len(text) > MAX_TOTAL_CHARS:
                log.warning(f"    STOP: Total content limit reached ({MAX_TOTAL_CHARS} chars)")
                break

            collected.append({
                "path": str(filepath),
                "type": content_type,
                "name": filepath.name,
                "content": text,
                "size": len(text),
            })
            total_chars += len(text)

            # Mark as processed
            processed[file_key] = file_mtime

    # Save updated state
    state["processed"] = processed
    state["last_run"] = datetime.now().isoformat()
    _save_digest_state(state_file, state)

    return collected


# --- Phase 2: LLM Analysis via GHCP SDK ---

async def run_digest(client: CopilotClient, config: dict):
    """Run a full digest cycle: collect content → analyze → write daily digest."""
    log.info("\n=== Digest cycle start ===")

    # Phase 1: Collect content
    log.info("Phase 1: Collecting content from input folders...")
    items = collect_content(config)

    if items:
        log.info(f"  Collected {len(items)} local items:")
        for item in items:
            log.info(f"    - [{item['type']}] {item['name']} ({item['size']} chars)")
    else:
        log.info("  No new local content.")

    # Phase 1b: Collect RSS feeds
    log.info("\nPhase 1b: Fetching RSS feeds...")
    articles = collect_feeds(config)
    if articles:
        log.info(f"  Collected {len(articles)} new articles")
    else:
        log.info("  No new articles.")

    # Phase 2: Send to GHCP SDK agent for analysis (always runs — WorkIQ queries happen here)
    log.info("Phase 2: Sending to agent for analysis + WorkIQ inbox scan...")

    async with agent_session(client, config, "digest", tools=get_tools()) as session:
        prompt = _build_digest_prompt(items, config, articles)
        log.info(f"  Prompt size: {len(prompt)} chars")
        log.info("  Agent working...")

        response = await session.send_and_wait({"prompt": prompt}, timeout=600)

        if not response:
            log.warning("No response from agent (timed out).")

    log.info("=== Digest cycle end ===")


def _load_previous_digest() -> str | None:
    """Load the most recent existing digest JSON to provide continuity."""
    digests_dir = OUTPUT_DIR / "digests"
    if not digests_dir.exists():
        return None
    json_files = sorted(digests_dir.glob("*.json"), reverse=True)
    if not json_files:
        return None
    try:
        data = json.loads(json_files[0].read_text(encoding="utf-8"))
        items = data.get("items", [])
        if not items:
            return None
        date = data.get("date", json_files[0].stem)
        lines = [f"## Previous Digest ({date}) — {len(items)} items were outstanding\n"]
        for item in items:
            priority = item.get("priority", "?")
            title = item.get("title", "?")
            item_type = item.get("type", "?")
            item_id = item.get("id", "?")
            lines.append(f"- [{priority.upper()}] **[{item_type}]** {title} (id: {item_id})")
        lines.append("\nCarry forward any items that are STILL outstanding. Drop items that have been resolved since then.")
        return "\n".join(lines)
    except Exception as e:
        log.warning(f"Could not load previous digest: {e}")
        return None


def _build_digest_prompt(items: list[dict], config: dict, articles: list[dict] | None = None) -> str:
    """Build the analysis prompt containing all collected content + RSS articles."""
    date_str = datetime.now().strftime("%Y-%m-%d")
    digest_cfg = config.get("digest", {})
    priorities = digest_cfg.get("priorities", [])
    priorities_str = "\n".join(f"- {p}" for p in priorities)

    # Group items by type
    by_type: dict[str, list[dict]] = {}
    for item in items:
        by_type.setdefault(item["type"], []).append(item)

    # Build content sections
    content_sections = []
    for content_type, type_items in by_type.items():
        section = f"### {content_type.title()} ({len(type_items)} files)\n"
        for item in type_items:
            section += f"\n---\n#### File: {item['name']}\n```\n{item['content']}\n```\n"
        content_sections.append(section)

    content_block = "\n".join(content_sections)

    # Build articles block from RSS feeds
    articles_block = ""
    if articles:
        article_lines = []
        for a in articles:
            article_lines.append(f"- [{a['source']}] **{a['title']}** ({a['published']})")
        articles_block = f"""
## Part C — External Intel ({len(articles)} articles from RSS feeds)
{chr(10).join(article_lines)}
"""

    # Load dismissed items
    actions = _load_actions()
    dismissed = actions.get("dismissed", [])
    notes = actions.get("notes", {})

    dismissed_block = ""
    if dismissed:
        dismissed_items = "\n".join(f"- {d['item']}" for d in dismissed)
        dismissed_block = f"""
## Previously Dismissed Items (DO NOT include these)
{dismissed_items}
"""

    notes_block = ""
    if notes:
        note_items = "\n".join(f"- **{k}**: {v['note']}" for k, v in notes.items())
        notes_block = f"""
## User Notes (context for your analysis)
{note_items}
"""

    # Load previous digest for continuity
    prev_digest = _load_previous_digest()
    prev_block = f"\n{prev_digest}\n" if prev_digest else ""

    # Load output rules from editable instruction file
    output_rules = _load_instruction("digest-output-rules", config)
    # Replace DATE placeholder with actual date
    output_rules = output_rules.replace("DATE", date_str)

    return f"""Generate a SHORT daily digest for {date_str}. This should be MAX 50 lines — only things I haven't dealt with yet.

## Your Priorities
{priorities_str}
{dismissed_block}{notes_block}{prev_block}
## Part A — Local Content (already collected)
{content_block}
{articles_block}
{output_rules}
"""


